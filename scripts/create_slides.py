#!/usr/bin/env python3
"""
Generador de slides completo para el Taller de Seguridad en Contraseñas
Crea slides duplicando la estructura de la plantilla
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from lxml import etree
import copy

def duplicate_slide(prs, index):
    """Duplica un slide existente"""
    template = prs.slides[index]
    slide_layout = template.slide_layout
    
    # Crear nuevo slide con el mismo layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    # Copiar las shapes del template al nuevo slide
    for shape in template.shapes:
        el = shape.element
        new_el = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    
    return new_slide

def update_slide_text(slide, replacements):
    """Actualiza textos en un slide según diccionario de reemplazos"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    original_text = run.text
                    for old, new in replacements.items():
                        if old in original_text:
                            run.text = original_text.replace(old, new)

def create_complete_presentation():
    """Crea la presentación completa del taller"""
    
    # Cargar plantilla
    prs = Presentation('slides/template.pptx')
    
    # Definición de contenido
    slides_content = [
        # Portada (ya existe - Slide 1)
        {
            "type": "cover",
            "replacements": {
                "Título de la charla o evento": "Taller Práctico de\nSeguridad en Contraseñas",
                "Fecha del evento": "Cátedras de Ciberseguridad 2025-2026"
            }
        },
        # Índice (Slide 2)
        {
            "type": "content",
            "title": "Contenido del Taller",
            "body": """1. Ataques Offline (Fuerza Bruta)
2. Ataques con Diccionario
3. Diccionarios Personalizados (CUPP)
4. Reglas de Mutación
5. Ataques Online (Hydra)
6. Credential Stuffing
7. Análisis de Stealers
8. Detección y Defensa (Blue Team)
9. Evasión y Rotación de IPs"""
        },
        # Módulo 1
        {
            "type": "content",
            "title": "Módulo 1: Ataques Offline - Fuerza Bruta",
            "body": """🎯 Objetivos:
• Comprender cómo se almacenan las contraseñas (hashes)
• Usar John the Ripper y Hashcat
• Evaluar complejidad temporal de ataques

🛠️ Herramientas: John the Ripper, Hashcat, Name-That-Hash

💡 Concepto clave:
   Hash = función one-way (irreversible)
   MD5, SHA-1 (débiles) vs bcrypt, Argon2 (fuertes)"""
        },
        # Módulo 2
        {
            "type": "content",
            "title": "Módulo 2: Ataques con Diccionario",
            "body": """📚 ¿Por qué funcionan los diccionarios?
• Las personas eligen contraseñas predecibles
• rockyou.txt: 14 millones de contraseñas reales

🔧 Recursos principales:
• SecLists (github.com/danielmiessler/SecLists)
• Weakpass (weakpass.com)
• Probable-Wordlists

⚡ Comando:
   hashcat -m 0 -a 0 hashes.txt rockyou.txt"""
        },
        # Módulo 3
        {
            "type": "content",
            "title": "Módulo 3: Diccionarios Personalizados",
            "body": """🎯 OSINT + Password Profiling

CUPP (Common User Passwords Profiler):
• Genera wordlists basadas en información personal
• Nombre, fechas, mascotas, equipos favoritos...

📊 Ejemplo de perfil:
   Nombre: Carlos García | Nacimiento: 15/03/1990
   Mascota: Luna | Equipo: Real Madrid

🔑 Genera: Carlos1990, Luna2024!, Garcia15#..."""
        },
        # Módulo 4
        {
            "type": "content",
            "title": "Módulo 4: Reglas de Mutación",
            "body": """🔄 Transformaciones automáticas:

password → Password, PASSWORD, p@ssw0rd,
           password123, password!, P@$$w0rd

📋 Reglas populares:
• best64.rule (básico, rápido)
• OneRuleToRuleThemAll (avanzado)
• Hob0Rules (políticas empresariales)

⚡ hashcat -m 0 hashes.txt wordlist.txt
           -r OneRuleToRuleThemAll.rule"""
        },
        # Módulo 5
        {
            "type": "content",
            "title": "Módulo 5: Ataques Online",
            "body": """🌐 Diferencia con Offline:
• Online: servicio activo, rate limiting, account lockout
• Más lento, detectable, pero acceso directo

🛠️ Herramientas:
• Hydra: SSH, FTP, HTTP, MySQL, SMB...
• FFUF: Web fuzzing moderno y rápido

⚡ Comando Hydra:
   hydra -l admin -P wordlist.txt
         ssh://192.168.1.1 -t 4"""
        },
        # Módulo 6
        {
            "type": "content",
            "title": "Módulo 6: Credential Stuffing",
            "body": """🔄 Ataque por reutilización de credenciales:

1. Filtración en Servicio A (LinkedIn 2012)
2. Usuario reutiliza password en Servicio B
3. Atacante prueba credenciales filtradas en B

📊 Estadísticas alarmantes:
• 65% de usuarios reutilizan contraseñas
• Have I Been Pwned: 13B+ credenciales filtradas

🛡️ Defensa: Passwords únicos + MFA + monitoring"""
        },
        # Módulo 7
        {
            "type": "content",
            "title": "Módulo 7: Análisis de Stealers",
            "body": """🦠 Info-Stealers (Redline, Raccoon, Vidar):
• Roban credenciales de navegadores
• Cookies de sesión (session hijacking)
• Wallets de criptomonedas

📋 MITRE ATT&CK Mapping:
• T1555.003: Credentials from Web Browsers
• T1539: Steal Web Session Cookie
• T1082: System Information Discovery

💰 Modelo: MaaS (Malware-as-a-Service)"""
        },
        # Módulo 8
        {
            "type": "content",
            "title": "Módulo 8: Detección y Defensa",
            "body": """🛡️ Perspectiva Blue Team:

Patrones en logs (grep, jq):
• Brute Force: misma IP, mismo user, muchos fallos
• Spraying: misma IP, muchos users, pocos fallos/user
• Stuffing: IPs rotatorias, algunos éxitos

🔐 Conditional Access (Zero Trust):
• Impossible Travel: Madrid 10:00 → Tokyo 11:00
• Device Compliance: ¿Antivirus activo?
• Risky IP: Tor, VPNs anónimas"""
        },
        # Módulo 9
        {
            "type": "content",
            "title": "Módulo 9: Evasión y Rotación de IPs",
            "body": """🚫 Rate Limiting y Bloqueos:
• WAFs bloquean IPs tras X intentos fallidos
• Error 429 Too Many Requests

🧅 Evasión con Tor / Proxy Pools:
• Enrutar tráfico vía Tor (SOCKS5)
• La IP de salida rota periódicamente
• Permite fuerza bruta lenta distribuida

🛠️ Práctica:
• Script Python + requests[socks]
• Contenedor Tor Proxy"""
        },
        # Herramientas
        {
            "type": "content",
            "title": "Resumen de Herramientas",
            "body": """🔓 Cracking Offline:
   John the Ripper, Hashcat, Name-That-Hash

📚 Wordlists:
   rockyou.txt, SecLists, CUPP, Pydictor

🌐 Ataques Online:
   Hydra, FFUF, Burp Suite

🦠 Análisis Malware:
   ANY.RUN, Joe Sandbox, YARA

🛡️ Defensa:
   grep, jq, Fail2Ban, EDR"""
        },
        # Entorno de práctica
        {
            "type": "content",
            "title": "Entorno de Práctica",
            "body": """🐳 Docker Compose incluido:

• attacker: Kali Linux con herramientas
• ssh-target: Servidor SSH vulnerable (puerto 2222)
• dvwa: Damn Vulnerable Web Application
• vulnerable-api: API Flask sin protección

⚡ Inicio rápido:
   docker-compose up -d
   docker-compose exec attacker /bin/bash

📁 github.com/edelahozuah/intro-to-password-workshop"""
        },
        # Ética
        {
            "type": "content",
            "title": "⚠️ Consideraciones Éticas",
            "body": """🚨 IMPORTANTE:

• NUNCA ejecutar técnicas sin autorización explícita
• El pentesting no autorizado es ILEGAL
• Usar SOLO el entorno Docker proporcionado
• Respetar leyes de privacidad y protección de datos

✅ Este taller es EXCLUSIVAMENTE educativo

🎓 Continúa aprendiendo:
   HackTheBox, TryHackMe, CTFs, OSCP"""
        },
        # Cierre
        {
            "type": "closing",
            "title": "¡Gracias!",
            "body": """🔗 Repositorio del Taller:
github.com/edelahozuah/intro-to-password-workshop

📧 Contacto:
Cátedras de Ciberseguridad UAH

#CátedrasCiber"""
        }
    ]
    
    # Procesar slides
    # Slide 1 (índice 0) - Portada
    slide1 = prs.slides[0]
    update_slide_text(slide1, slides_content[0]["replacements"])
    
    # Slide 2 (índice 1) - Actualizar como índice
    slide2 = prs.slides[1]
    update_slide_text(slide2, {
        "Título de la charla o evento": slides_content[1]["title"],
        "Fecha del evento": "8 módulos prácticos"
    })
    
    # Slides 3 y 4 ya tienen la estructura de contenido
    for i, content_slide in enumerate([prs.slides[2], prs.slides[3]]):
        content_idx = i + 2  # slides_content[2] y [3]
        if content_idx < len(slides_content):
            content = slides_content[content_idx]
            for shape in content_slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if "Título 1" in text:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                run.text = content["title"]
                    elif "Cuerpo de texto" in text:
                        # Limpiar y añadir nuevo contenido
                        tf = shape.text_frame
                        tf.clear()
                        lines = content["body"].split('\n')
                        for j, line in enumerate(lines):
                            if j == 0:
                                p = tf.paragraphs[0]
                            else:
                                p = tf.add_paragraph()
                            run = p.add_run()
                            run.text = line
    
    # Duplicar slide 3 (índice 2) para crear más slides de contenido
    template_idx = 2  # Usamos slide 3 como template (tiene título + cuerpo)
    
    # Crear slides para los módulos 3-8 y adicionales
    for content_idx in range(4, len(slides_content) - 1):  # Saltamos portada, índice, y los 2 primeros módulos
        content = slides_content[content_idx]
        new_slide = duplicate_slide(prs, template_idx)
        
        # Actualizar el contenido del nuevo slide
        for shape in new_slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if "Título 1" in text or "Módulo" in text:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = content["title"]
                elif "Cuerpo" in text or len(text) > 50:
                    tf = shape.text_frame
                    tf.clear()
                    lines = content["body"].split('\n')
                    for j, line in enumerate(lines):
                        if j == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        run = p.add_run()
                        run.text = line
    
    # Actualizar slide de cierre (último existente - slide 6)
    closing_content = slides_content[-1]
    slide6 = prs.slides[5]
    update_slide_text(slide6, {
        "Título de la charla o evento": closing_content["title"]
    })
    
    # Guardar
    output_path = 'slides/Taller_Seguridad_Contraseñas.pptx'
    prs.save(output_path)
    print(f"✅ Presentación guardada: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    
    return output_path

if __name__ == "__main__":
    create_complete_presentation()
