#!/usr/bin/env python3
"""
Script para explorar la estructura de una plantilla PPTX
"""
from pptx import Presentation
from pptx.util import Inches, Pt

# Cargar plantilla
prs = Presentation('slides/template.pptx')

print("=" * 60)
print("ESTRUCTURA DE LA PLANTILLA")
print("=" * 60)

# Tamaño de slide
print(f"\nTamaño de slide: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")

# Layouts disponibles
print(f"\nLayouts disponibles ({len(prs.slide_layouts)}):")
for i, layout in enumerate(prs.slide_layouts):
    print(f"  [{i}] {layout.name}")

# Slides existentes
print(f"\nSlides en la plantilla ({len(prs.slides)}):")
for i, slide in enumerate(prs.slides):
    layout_name = slide.slide_layout.name
    print(f"\n  Slide {i+1} - Layout: '{layout_name}'")
    
    # Shapes en el slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text[:50].replace('\n', ' ')
            print(f"    - Shape: {shape.shape_type}, Text: '{text}...'")
        else:
            print(f"    - Shape: {shape.shape_type}")

print("\n" + "=" * 60)
